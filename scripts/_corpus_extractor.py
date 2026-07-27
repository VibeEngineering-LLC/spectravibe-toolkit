from __future__ import annotations
import sys
try:
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
except Exception:
    pass
"""Parallel full-text extractor for ALL reference documents.

Walks references/books/, extracts text from every .pdf/.docx/.pptx into
references/_extracted_corpus/{relative_path}.md, using ProcessPoolExecutor
across all CPU cores. Designed for a one-shot complete corpus dump so that
downstream sub-agents can do deep reads from local files (no per-page PDF
parsing in the agent context).

Run:
    python scripts/_corpus_extractor.py
"""

import json
import os
import sys
import time
import traceback
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
BOOKS = ROOT / "references" / "books"
OUT = ROOT / "references" / "_extracted_corpus"


def extract_pdf(src: Path) -> tuple[str, dict]:
    import fitz  # pymupdf

    doc = fitz.open(str(src))
    parts: list[str] = []
    meta = {
        "n_pages": doc.page_count,
        "n_chars": 0,
        "n_images": 0,
        "has_text_layer": False,
    }
    for i, page in enumerate(doc):
        text = page.get_text("text") or ""
        if text.strip():
            meta["has_text_layer"] = True
        try:
            img_list = page.get_images(full=True)
            meta["n_images"] += len(img_list)
        except Exception:
            pass
        parts.append(f"\n\n## PAGE {i+1}\n\n{text}")
        meta["n_chars"] += len(text)
    doc.close()
    return "".join(parts), meta


def extract_docx(src: Path) -> tuple[str, dict]:
    import docx

    d = docx.Document(str(src))
    parts: list[str] = []
    n_tbl = 0
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            sty = p.style.name if p.style else ""
            if sty.startswith("Heading"):
                lvl = "".join(ch for ch in sty if ch.isdigit()) or "1"
                parts.append(f"\n\n{'#' * (int(lvl) + 1)} {t}\n")
            else:
                parts.append(t)
    for tbl in d.tables:
        n_tbl += 1
        parts.append(f"\n\n### Table {n_tbl}\n")
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    text = "\n".join(parts)
    return text, {"n_chars": len(text), "n_tables": n_tbl, "n_paragraphs": len(d.paragraphs)}


def extract_pptx(src: Path) -> tuple[str, dict]:
    import pptx

    p = pptx.Presentation(str(src))
    parts: list[str] = []
    n_imgs = 0
    n_tables = 0
    for i, slide in enumerate(p.slides, 1):
        parts.append(f"\n\n## SLIDE {i}\n")
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            parts.append(t)
                if getattr(shape, "has_table", False):
                    n_tables += 1
                    parts.append(f"\n### Table {n_tables}\n")
                    tbl = shape.table
                    for row in tbl.rows:
                        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                        parts.append("| " + " | ".join(cells) + " |")
                if shape.shape_type == 13:  # picture
                    n_imgs += 1
            except Exception:
                continue
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        if notes:
            parts.append(f"\n[NOTES]: {notes}")
    text = "\n".join(parts)
    return text, {"n_chars": len(text), "n_slides": len(p.slides), "n_images": n_imgs, "n_tables": n_tables}


def extract_doc(src: Path) -> tuple[str, dict]:
    """Old .doc — try LibreOffice headless conversion via soffice if present."""
    # Skip; we don't have a soffice install guarantee. Return marker.
    return "[OLD .doc — NOT EXTRACTED, convert manually]", {"skipped": True}


def extract_ppt(src: Path) -> tuple[str, dict]:
    return "[OLD .ppt — NOT EXTRACTED, convert manually]", {"skipped": True}


HANDLERS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".doc": extract_doc,
    ".ppt": extract_ppt,
}


def _safe_relpath(src: Path) -> Path:
    return src.relative_to(BOOKS)


def process_one(src_str: str) -> dict:
    src = Path(src_str)
    ext = src.suffix.lower()
    handler = HANDLERS.get(ext)
    if not handler:
        return {"src": str(src), "skipped": True, "reason": "unsupported ext"}
    rel = _safe_relpath(src)
    out_md = OUT / rel.with_suffix(rel.suffix + ".md")
    out_md.parent.mkdir(parents=True, exist_ok=True)
    t0 = time.time()
    try:
        text, meta = handler(src)
    except Exception as e:
        return {
            "src": str(src),
            "error": f"{type(e).__name__}: {e}",
            "trace": traceback.format_exc()[-400:],
        }
    out_md.write_text(
        f"<!-- src: {rel} -->\n<!-- ext: {ext} -->\n<!-- meta: {json.dumps(meta, ensure_ascii=False)} -->\n\n# {rel.name}\n\n{text}",
        encoding="utf-8",
    )
    return {
        "src": str(rel),
        "out": str(out_md.relative_to(ROOT)),
        "elapsed_s": round(time.time() - t0, 2),
        **meta,
    }


def discover() -> list[Path]:
    files: list[Path] = []
    for p in BOOKS.rglob("*"):
        if p.is_file() and p.suffix.lower() in HANDLERS:
            files.append(p)
    return files


def main() -> int:
    files = discover()
    print(f"Discovered {len(files)} documents in {BOOKS}", flush=True)
    OUT.mkdir(parents=True, exist_ok=True)
    workers = max(1, (os.cpu_count() or 4) - 1)
    print(f"Using {workers} workers", flush=True)

    results: list[dict] = []
    t_start = time.time()
    with ProcessPoolExecutor(max_workers=workers) as ex:
        futs = {ex.submit(process_one, str(f)): f for f in files}
        done = 0
        for fut in as_completed(futs):
            r = fut.result()
            results.append(r)
            done += 1
            tag = "ERR" if "error" in r else ("SKIP" if r.get("skipped") else "OK ")
            print(f"[{done:3d}/{len(files)}] {tag} {r.get('src','?')} ({r.get('elapsed_s','?')}s)", flush=True)

    elapsed = round(time.time() - t_start, 2)
    summary = {
        "n_files": len(files),
        "n_ok": sum(1 for r in results if "error" not in r and not r.get("skipped")),
        "n_err": sum(1 for r in results if "error" in r),
        "n_skip": sum(1 for r in results if r.get("skipped")),
        "total_chars": sum(r.get("n_chars", 0) for r in results),
        "total_pages": sum(r.get("n_pages", 0) for r in results),
        "total_slides": sum(r.get("n_slides", 0) for r in results),
        "elapsed_s": elapsed,
        "results": results,
    }
    (OUT / "_extraction_report.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(
        f"\nDONE in {elapsed}s | ok={summary['n_ok']} err={summary['n_err']} skip={summary['n_skip']} "
        f"| {summary['total_chars']:,} chars across {summary['total_pages']} PDF pages + {summary['total_slides']} slides"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
